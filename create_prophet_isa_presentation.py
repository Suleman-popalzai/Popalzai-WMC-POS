#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Professional PowerPoint Presentation about Prophet Isa (AS)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme (Professional Islamic theme)
    PRIMARY_COLOR = RGBColor(0, 102, 102)  # Teal
    SECONDARY_COLOR = RGBColor(255, 255, 255)  # White
    ACCENT_COLOR = RGBColor(204, 153, 51)  # Gold
    TEXT_COLOR = RGBColor(51, 51, 51)  # Dark Gray
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Prophet Isa (Jesus) عليه السلام"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = SECONDARY_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Eine umfassende Betrachtung aus islamischer Perspektive"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = ACCENT_COLOR
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Inhaltsverzeichnis
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide2.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = SECONDARY_COLOR
    
    # Header
    add_header(slide2, "Inhaltsverzeichnis", PRIMARY_COLOR, SECONDARY_COLOR)
    
    # Content
    content_box = slide2.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    topics = [
        "1. Einleitung und Bedeutung",
        "2. Geburt und Wunder",
        "3. Prophetentum und Botschaft",
        "4. Wunder und Zeichen",
        "5. Die Kreuzigung im Islam",
        "6. Isa (AS) im Quran",
        "7. Rückkehr und Endzeit",
        "8. Zusammenfassung"
    ]
    
    for i, topic in enumerate(topics):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = topic
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(12)
        p.level = 0
    
    # Slide 3: Einleitung und Bedeutung
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide3.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = SECONDARY_COLOR
    
    add_header(slide3, "Einleitung und Bedeutung", PRIMARY_COLOR, SECONDARY_COLOR)
    
    content_box = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    points = [
        "• Prophet Isa (AS) ist einer der bedeutendsten Propheten im Islam",
        "• Im Quran wird er 25 Mal namentlich erwähnt",
        "• Er wird als 'Masih' (Messias) bezeichnet",
        "• Einer der fünf Ulul-Azm Propheten (Propheten mit großer Entschlossenheit)",
        "• Seine Mutter Maryam (Maria) ist die einzige Frau, die im Quran namentlich erwähnt wird",
        "• Eine ganze Sure (Sure 19) ist nach seiner Mutter benannt"
    ]
    
    for i, point in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(10)
    
    # Slide 4: Geburt und Wunder
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide4.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = SECONDARY_COLOR
    
    add_header(slide4, "Die wundersame Geburt", PRIMARY_COLOR, SECONDARY_COLOR)
    
    content_box = slide4.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    points = [
        "• Maryam (AS) empfing Isa durch Allahs Wort 'Kun' (Sei!)",
        "• Geburt ohne Vater - ein Zeichen Allahs Allmacht",
        "• Maryam zog sich während der Schwangerschaft zurück",
        "• Geburt unter einer Dattelpalme",
        "• Das Wunder: Isa (AS) sprach als Neugeborenes",
        "• Er verteidigte die Ehre seiner Mutter",
        "",
        "Quran (19:30-31): 'Er sprach: Ich bin wahrlich Allahs Diener; Er hat mir die Schrift gegeben und mich zu einem Propheten gemacht.'"
    ]
    
    for i, point in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        if point.startswith("Quran"):
            p.font.size = Pt(15)
            p.font.italic = True
            p.font.color.rgb = ACCENT_COLOR
        else:
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)
    
    # Slide 5: Prophetentum und Botschaft
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide5.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = SECONDARY_COLOR
    
    add_header(slide5, "Prophetentum und Botschaft", PRIMARY_COLOR, SECONDARY_COLOR)
    
    content_box = slide5.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    points = [
        "Kernbotschaften:",
        "• Tawhid - Der Glaube an den Einen Gott",
        "• Aufruf zur Anbetung Allahs allein",
        "• Bestätigung der Tora und Bringung des Injil (Evangelium)",
        "• Verkündigung ethischer und moralischer Werte",
        "• Gerechtigkeit, Barmherzigkeit und Mitgefühl",
        "",
        "Quran (3:49-50): 'Und als Gesandter zu den Kindern Israels: Ich bin mit einem Zeichen von eurem Herrn zu euch gekommen...'"
    ]
    
    for i, point in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        if point == "Kernbotschaften:":
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_COLOR
        elif point.startswith("Quran"):
            p.font.size = Pt(15)
            p.font.italic = True
            p.font.color.rgb = ACCENT_COLOR
        else:
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)
    
    # Slide 6: Wunder und Zeichen
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide6.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = SECONDARY_COLOR
    
    add_header(slide6, "Wunder und Zeichen (Mu'jizat)", PRIMARY_COLOR, SECONDARY_COLOR)
    
    content_box = slide6.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    points = [
        "Allah gab Isa (AS) außergewöhnliche Wunder:",
        "",
        "• Heilung von Blinden und Aussätzigen",
        "• Erweckung der Toten mit Allahs Erlaubnis",
        "• Erschaffung eines Vogels aus Lehm, der lebendig wurde",
        "• Wissen über verborgene Dinge",
        "• Herabsendung eines Tisches vom Himmel (Ma'ida)",
        "",
        "Wichtig: Alle Wunder geschahen durch Allahs Willen und Erlaubnis, nicht durch eigene Macht"
    ]
    
    for i, point in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        if i == 0:
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_COLOR
        elif point.startswith("Wichtig:"):
            p.font.size = Pt(16)
            p.font.italic = True
            p.font.color.rgb = ACCENT_COLOR
        else:
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)
    
    # Slide 7: Die Kreuzigung im Islam
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide7.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = SECONDARY_COLOR
    
    add_header(slide7, "Die Kreuzigung im Islam", PRIMARY_COLOR, SECONDARY_COLOR)
    
    content_box = slide7.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    points = [
        "Islamische Sichtweise:",
        "",
        "• Isa (AS) wurde NICHT gekreuzigt",
        "• Allah rettete ihn vor seinen Feinden",
        "• Jemand anderes wurde ihm ähnlich gemacht",
        "• Er wurde zu Allah erhoben (Himmelfahrt)",
        "",
        "Quran (4:157-158): 'Sie haben ihn aber nicht getötet und nicht gekreuzigt, sondern es erschien ihnen so... Vielmehr hat Allah ihn zu Sich erhoben.'",
        "",
        "Dies ist ein fundamentaler Unterschied zum christlichen Glauben"
    ]
    
    for i, point in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        if i == 0:
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_COLOR
        elif point.startswith("Quran"):
            p.font.size = Pt(15)
            p.font.italic = True
            p.font.color.rgb = ACCENT_COLOR
        elif point.startswith("Dies ist"):
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_COLOR
        else:
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)
    
    # Slide 8: Isa (AS) im Quran
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide8.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = SECONDARY_COLOR
    
    add_header(slide8, "Isa (AS) im Heiligen Quran", PRIMARY_COLOR, SECONDARY_COLOR)
    
    content_box = slide8.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    points = [
        "Titel und Bezeichnungen im Quran:",
        "",
        "• Al-Masih (Der Messias)",
        "• Ibn Maryam (Sohn der Maria)",
        "• Rasulullah (Gesandter Allahs)",
        "• Kalimatullah (Wort Allahs)",
        "• Ruhun minhu (Geist von Ihm)",
        "• Wajihan fid-dunya wal-akhira (Angesehen in dieser Welt und im Jenseits)",
        "",
        "Wichtige Suren: Al-Imran (3), An-Nisa (4), Al-Ma'ida (5), Maryam (19)"
    ]
    
    for i, point in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        if i == 0 or point.startswith("Wichtige"):
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_COLOR
        else:
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)
    
    # Slide 9: Rückkehr und Endzeit
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide9.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = SECONDARY_COLOR
    
    add_header(slide9, "Die Rückkehr von Isa (AS)", PRIMARY_COLOR, SECONDARY_COLOR)
    
    content_box = slide9.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    points = [
        "Zeichen der Endzeit:",
        "",
        "• Isa (AS) wird vor dem Jüngsten Tag zur Erde zurückkehren",
        "• Er wird als gerechter Richter erscheinen",
        "• Er wird den Dajjal (Antichristen) besiegen",
        "• Er wird nach der Scharia des Propheten Muhammad ﷺ urteilen",
        "• Er wird heiraten und Kinder haben",
        "• Er wird sterben und begraben werden",
        "",
        "Seine Rückkehr ist ein großes Zeichen für das Ende der Zeit"
    ]
    
    for i, point in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        if i == 0:
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_COLOR
        elif point.startswith("Seine Rückkehr"):
            p.font.size = Pt(16)
            p.font.italic = True
            p.font.color.rgb = ACCENT_COLOR
        else:
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)
    
    # Slide 10: Zusammenfassung
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide10.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = SECONDARY_COLOR
    
    add_header(slide10, "Zusammenfassung", PRIMARY_COLOR, SECONDARY_COLOR)
    
    content_box = slide10.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    points = [
        "Kernpunkte:",
        "",
        "✓ Isa (AS) ist ein hochverehrter Prophet im Islam",
        "✓ Seine wundersame Geburt zeigt Allahs Allmacht",
        "✓ Er brachte die Botschaft des Tawhid (Monotheismus)",
        "✓ Er wurde nicht gekreuzigt, sondern zu Allah erhoben",
        "✓ Er wird am Ende der Zeit zurückkehren",
        "✓ Muslime lieben und respektieren ihn zutiefst",
        "",
        "Der Glaube an alle Propheten, einschließlich Isa (AS), ist ein fundamentaler Bestandteil des islamischen Glaubens"
    ]
    
    for i, point in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        if i == 0:
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_COLOR
        elif point.startswith("Der Glaube"):
            p.font.size = Pt(16)
            p.font.italic = True
            p.font.color.rgb = ACCENT_COLOR
        else:
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(10)
    
    # Slide 11: Abschluss
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide11.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Thank you message
    thank_box = slide11.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    thank_frame = thank_box.text_frame
    thank_frame.text = "Vielen Dank für Ihre Aufmerksamkeit"
    thank_para = thank_frame.paragraphs[0]
    thank_para.font.size = Pt(44)
    thank_para.font.bold = True
    thank_para.font.color.rgb = SECONDARY_COLOR
    thank_para.alignment = PP_ALIGN.CENTER
    
    # Arabic text
    arabic_box = slide11.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    arabic_frame = arabic_box.text_frame
    arabic_frame.text = "صلى الله عليه وسلم"
    arabic_para = arabic_frame.paragraphs[0]
    arabic_para.font.size = Pt(32)
    arabic_para.font.color.rgb = ACCENT_COLOR
    arabic_para.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save('/vercel/sandbox/Prophet_Isa_AS_Presentation.pptx')
    print("✓ Präsentation erfolgreich erstellt: Prophet_Isa_AS_Presentation.pptx")

def add_header(slide, title_text, bg_color, text_color):
    """Add a professional header to a slide"""
    # Header background
    header_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(1.3)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = bg_color
    header_shape.line.color.rgb = bg_color
    
    # Header text
    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = header_text.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT

if __name__ == "__main__":
    create_presentation()
